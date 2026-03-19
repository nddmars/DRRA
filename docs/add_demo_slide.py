"""
Add WSG Demo Slide to RSA 2026 Presentation.

Inserts slide after slide 15 (WSG Framework image slide) showing
Resilience Forge (DRRA) as the live open-source implementation of WSG.

Run:
    python docs/add_demo_slide.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import sys

PPTX_PATH = r"C:\surya\OneDrive\EB1A\final\_misc\Conferences-speaker\industry\2026-RSA\IMT-T09 - Beyond Prevention\u2014Architecting Cyber-Resilient Defense Against Ransomware v4.pptx"
OUT_PATH  = r"C:\surya\OneDrive\EB1A\final\_misc\Conferences-speaker\industry\2026-RSA\IMT-T09 - Beyond Prevention\u2014Architecting Cyber-Resilient Defense Against Ransomware v5-demo.pptx"

# ── Palette (matches "Midnight Executive" / dark cyber theme) ──────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
WALL_C    = RGBColor(0x00, 0x7B, 0xFF)   # electric blue
SQUAT_C   = RGBColor(0xFF, 0xA5, 0x00)   # amber
GRAB_C    = RGBColor(0x00, 0xC8, 0x6E)   # green
FEED_C    = RGBColor(0xAA, 0x55, 0xFF)   # purple
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
DARK_CARD = RGBColor(0x1A, 0x2E, 0x44)
ACCENT    = RGBColor(0xFF, 0x45, 0x4F)   # red accent


def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_label_box(slide, label, sublabel, l, t, w, h, color):
    """Colored card with a large label and smaller sublabel."""
    add_rect(slide, l, t, w, h, DARK_CARD)
    # left accent bar
    add_rect(slide, l, t, 0.04, h, color)
    # label
    add_text(slide, label, l + 0.12, t + 0.08, w - 0.16, 0.35,
             font_size=15, bold=True, color=color)
    # sublabel (multiline)
    add_text(slide, sublabel, l + 0.12, t + 0.40, w - 0.16, h - 0.48,
             font_size=9.5, bold=False, color=LIGHT)


def add_arrow(slide, l, t, w=0.01, h=0.35, vertical=True):
    """Simple thin rectangle as arrow."""
    add_rect(slide, l, t, 0.04 if vertical else w, h if vertical else 0.04,
             RGBColor(0x55, 0x77, 0x99))


def build_slide(prs: Presentation) -> None:
    """Append the demo slide to the presentation."""
    # Use a blank slide layout
    blank_layout = prs.slide_layouts[6]  # "Blank"
    slide = prs.slides.add_slide(blank_layout)

    W, H = 10.0, 7.5  # standard widescreen

    # ── Background ─────────────────────────────────────────────────────────
    bg = add_rect(slide, 0, 0, W, H, BG)

    # ── Slide title ─────────────────────────────────────────────────────────
    add_text(
        slide,
        "LIVE REFERENCE ARCHITECTURE: Resilience Forge (DRRA)",
        0.3, 0.12, 9.4, 0.45,
        font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )
    add_text(
        slide,
        "Open-source WSG proof-of-concept  |  github.com/resilience-forge",
        0.3, 0.54, 9.4, 0.28,
        font_size=11, bold=False, color=LIGHT, align=PP_ALIGN.LEFT,
    )
    # thin separator line
    add_rect(slide, 0.3, 0.82, 9.4, 0.025, RGBColor(0x33, 0x55, 0x77))

    # ── Top row: FORGE banner ────────────────────────────────────────────────
    add_rect(slide, 0.3, 0.92, 9.4, 0.44, RGBColor(0x1A, 0x2E, 0x44))
    add_rect(slide, 0.3, 0.92, 0.06, 0.44, RGBColor(0x88, 0x88, 0x88))
    add_text(slide, "THE FORGE  (Simulation Engine)", 0.5, 0.94, 3.5, 0.22,
             font_size=11, bold=True, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text(slide,
             "HoneypotGenerator  |  PayloadSimulator  |  ResiliencePayloadGenerator  |  KerberosTest",
             0.5, 1.14, 9.0, 0.20, font_size=9, color=RGBColor(0x88, 0x99, 0xAA))

    # ── Down arrow from FORGE ────────────────────────────────────────────────
    add_arrow(slide, 4.96, 1.38, vertical=True, h=0.28)

    # ── Three pillar cards ───────────────────────────────────────────────────
    CARD_T = 1.70
    CARD_H = 1.90
    GAP    = 0.15

    # WALL card
    add_label_box(
        slide,
        "WALL  (PREVENT / DETECT)",
        "sentinel_service.py  |  detector.py\n"
        "watchers/lib.rs (Rust file events)\n"
        "sigma/ransomware_mass_modification.yml\n"
        "\nEntropy > 0.85 => encryption detected\n"
        "Mass-mod > 15% / 60s => critical alert\n"
        "Gemini LLM => plain-English insight",
        0.3, CARD_T, 3.05, CARD_H, WALL_C,
    )

    # SQUAT card
    add_label_box(
        slide,
        "SQUAT  (SURVIVE / CONTAIN)",
        "shield_service.py  |  MicroSegmentationService\n"
        "opa/blast_radius_containment.rego\n"
        "opa/isolation_vlan_rules.rego\n"
        "\n>95% confidence => auto-isolate (VLAN 9999)\n"
        "70-95% => isolate + approval gate\n"
        "MTTC target: <60s  |  Achieved: 2.1s",
        3.65, CARD_T, 3.05, CARD_H, SQUAT_C,
    )

    # GRAB card
    add_label_box(
        slide,
        "GRAB  (CONTROL / RECOVER)",
        "shield_service.py  |  RecoveryOrchestrator\n"
        "MinIO Object Lock (compliance mode)\n"
        "playbooks/active_ransomware_response.yml\n"
        "\nForensic evidence => immutable (90-day lock)\n"
        "Snapshot restore => 8 parallel threads\n"
        "Data loss target: <0.1%",
        7.00, CARD_T, 2.70, CARD_H, GRAB_C,
    )

    # ── Horizontal connector arrows between cards ────────────────────────────
    ARR_Y = CARD_T + CARD_H / 2 - 0.02
    add_rect(slide, 3.35, ARR_Y, 0.30, 0.04, SQUAT_C)   # WALL -> SQUAT
    add_rect(slide, 6.70, ARR_Y, 0.30, 0.04, GRAB_C)    # SQUAT -> GRAB

    # ── Feedback Loop band ───────────────────────────────────────────────────
    FEED_T = CARD_T + CARD_H + 0.18
    add_rect(slide, 0.3, FEED_T, 9.4, 0.55, RGBColor(0x18, 0x0F, 0x2E))
    add_rect(slide, 0.3, FEED_T, 0.06, 0.55, FEED_C)
    add_text(slide, "FEEDBACK LOOP  (Closed-Loop Learning)",
             0.5, FEED_T + 0.04, 3.8, 0.25,
             font_size=11, bold=True, color=FEED_C)
    add_text(slide,
             "feedback_service.py  |  DefensibilityScorer  |  ThresholdAdaptor  |  "
             "DevSecOps CI/CD (test_gauntlet.py)  |  vector.toml => MinIO immutable telemetry",
             0.5, FEED_T + 0.27, 9.0, 0.22,
             font_size=9, color=LIGHT)

    # ── Bottom metrics strip ─────────────────────────────────────────────────
    MET_T = FEED_T + 0.62
    metrics = [
        ("DI Score", "49=>94", "Improving\nacross 5 incidents", FEED_C),
        ("WALL MTTD", "<3.2 min", "vs 241-day\nindustry avg", WALL_C),
        ("SQUAT MTTC", "2.1 sec", "vs 48-min\nbreakout speed", SQUAT_C),
        ("GRAB Recovery", ">99.9%", "Data recovered\nfrom immutable snap", GRAB_C),
        ("FP Adaptation", "0.87", "Entropy threshold\nauto-tuned from 0.85", ACCENT),
    ]
    BOX_W = 1.82
    for i, (label, value, sub, color) in enumerate(metrics):
        bx = 0.3 + i * (BOX_W + 0.06)
        add_rect(slide, bx, MET_T, BOX_W, 1.0, DARK_CARD)
        add_rect(slide, bx, MET_T, BOX_W, 0.04, color)          # top accent
        add_text(slide, label, bx + 0.08, MET_T + 0.08, BOX_W - 0.16, 0.22,
                 font_size=8.5, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, value, bx + 0.08, MET_T + 0.27, BOX_W - 0.16, 0.35,
                 font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, sub, bx + 0.08, MET_T + 0.60, BOX_W - 0.16, 0.36,
                 font_size=8, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)

    # ── Footer ───────────────────────────────────────────────────────────────
    add_text(slide,
             "Stack: Python (FastAPI)  |  Rust  |  MinIO  |  OPA (Rego)  |  Sigma  |  "
             "Kafka  |  React  |  Ansible  |  Gemini 2.5 Flash  |  Docker",
             0.3, 7.18, 9.4, 0.25,
             font_size=7.5, color=RGBColor(0x55, 0x77, 0x99), align=PP_ALIGN.CENTER)


def main():
    import os
    pptx_path = PPTX_PATH
    out_path  = OUT_PATH

    # Fallback if Unicode escape didn't resolve (Windows path)
    if not os.path.exists(pptx_path):
        pptx_path = pptx_path.replace("\\u2014", "\u2014")
        out_path  = out_path.replace("\\u2014", "\u2014")

    if not os.path.exists(pptx_path):
        print(f"ERROR: Source file not found:\n  {pptx_path}")
        sys.exit(1)

    print(f"Loading: {pptx_path}")
    prs = Presentation(pptx_path)
    print(f"  Existing slides: {len(prs.slides)}")

    build_slide(prs)
    print(f"  Slides after insert: {len(prs.slides)}")

    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
