"""
DRRA Conference Demo - PowerPoint Generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x8B, 0xFF)   # Bright blue
ACCENT_RED  = RGBColor(0xFF, 0x3B, 0x3B)   # Alert red
ACCENT_GRN  = RGBColor(0x00, 0xE6, 0x76)   # Success green
ACCENT_YLW  = RGBColor(0xFF, 0xC3, 0x00)   # Warning yellow
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB0, 0xBE, 0xC5)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # fully blank


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h, size=Pt(18), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def accent_bar(slide, color=ACCENT_BLUE):
    box(slide, 0, 0, 13.33, 0.08, fill_color=color)


def slide_number(slide, n):
    txt(slide, str(n), 12.8, 7.1, 0.5, 0.3, size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
box(s, 0, 0.08, 13.33, 7.42, fill_color=DARK_BG)

# Big title
txt(s, "DRRA", 1, 1.5, 11, 1.8, size=Pt(96), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
txt(s, "Distributed Resilience & Recovery Architecture", 1, 3.1, 11, 0.7,
    size=Pt(26), bold=False, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Ransomware Defense — Live Demo", 1, 3.8, 11, 0.5,
    size=Pt(18), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Tag line
box(s, 3.5, 5.0, 6.3, 0.6, fill_color=ACCENT_BLUE)
txt(s, "  Detect  •  Isolate  •  Recover  •  Prove It  ", 3.5, 5.0, 6.3, 0.6,
    size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(s, 1)

# ──────────────────────────────────────────────
# SLIDE 2 — The Problem
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, ACCENT_RED)

txt(s, "The Ransomware Problem", 0.5, 0.2, 12, 0.6, size=Pt(32), bold=True, color=WHITE)

stats = [
    ("$4.9M",   "Average ransom\npayment (2024)"),
    ("21 Days",  "Average downtime\nafter attack"),
    ("83%",      "Orgs attacked\nmore than once"),
    ("< 8 sec",  "DRRA detection\ntime"),
]
colors = [ACCENT_RED, ACCENT_RED, ACCENT_RED, ACCENT_GRN]

for i, (val, label) in enumerate(stats):
    x = 0.4 + i * 3.15
    box(s, x, 1.1, 2.8, 2.2, fill_color=RGBColor(0x1A, 0x2B, 0x3C), line_color=colors[i], line_width=Pt(2))
    txt(s, val,   x, 1.2, 2.8, 1.0, size=Pt(38), bold=True, color=colors[i], align=PP_ALIGN.CENTER)
    txt(s, label, x, 2.2, 2.8, 0.9, size=Pt(14), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "Traditional defenses focus on prevention. DRRA focuses on defensibility — how fast you detect, contain, and recover.",
    0.5, 3.6, 12.3, 0.8, size=Pt(17), color=WHITE)

txt(s, "Prevention fails. Defensibility wins.", 0.5, 4.5, 12.3, 0.6,
    size=Pt(22), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

slide_number(s, 2)

# ──────────────────────────────────────────────
# SLIDE 3 — Architecture Overview
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)

txt(s, "DRRA Architecture", 0.5, 0.2, 12, 0.6, size=Pt(32), bold=True, color=WHITE)

components = [
    ("FORGE",   "Attack\nSimulation",   ACCENT_RED,  0.4),
    ("WATCHER", "File System\nMonitor", ACCENT_YLW,  3.0),
    ("VIGIL",   "Behavioral\nDetection",ACCENT_BLUE, 5.6),
    ("SHIELD",  "Isolation &\nRecovery",ACCENT_GRN,  8.2),
]

for name, desc, color, x in components:
    box(s, x, 1.1, 2.3, 2.0, fill_color=RGBColor(0x1A, 0x2B, 0x3C), line_color=color, line_width=Pt(2))
    txt(s, name, x, 1.2, 2.3, 0.7, size=Pt(20), bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc, x, 1.9, 2.3, 0.9, size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Arrows between boxes
for ax in [2.75, 5.35, 7.95]:
    txt(s, "→", ax, 1.7, 0.5, 0.5, size=Pt(28), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Infrastructure row
infra = [
    ("Kafka",      "Event Streaming"),
    ("PostgreSQL", "Audit Logs"),
    ("MinIO",      "Immutable Storage"),
    ("Redis",      "Caching"),
    ("Prometheus", "Metrics"),
    ("Grafana",    "Dashboards"),
]
txt(s, "Infrastructure Layer", 0.4, 3.4, 12, 0.4, size=Pt(13), color=LIGHT_GRAY)
for i, (name, desc) in enumerate(infra):
    x = 0.4 + i * 2.1
    box(s, x, 3.8, 1.9, 1.1, fill_color=RGBColor(0x12, 0x22, 0x30), line_color=LIGHT_GRAY, line_width=Pt(1))
    txt(s, name, x, 3.85, 1.9, 0.45, size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, x, 4.3, 1.9, 0.4,  size=Pt(10), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "All services containerised with Docker Compose  •  REST API via FastAPI  •  Rust file watcher for low-latency detection",
    0.4, 5.15, 12.5, 0.5, size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 3)

# ──────────────────────────────────────────────
# SLIDE 4 — Live Demo Flow
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, ACCENT_YLW)

txt(s, "Live Demo Flow", 0.5, 0.2, 12, 0.6, size=Pt(32), bold=True, color=WHITE)

steps = [
    ("1", "FORGE Deploy",          "POST /api/v1/forge/deploy",           ACCENT_RED,  "Simulate honeypot ransomware attack"),
    ("2", "VIGIL Detection",       "GET  /api/v1/vigil/events?payload_id=","ACCENT_BLUE","3 signatures: entropy · VSS · lateral movement"),
    ("3", "SHIELD Isolate",        "POST /api/v1/shield/isolate",         ACCENT_GRN,  "VLAN micro-segmentation in seconds"),
    ("4", "SHIELD Recovery",       "POST /api/v1/shield/recovery/create", ACCENT_GRN,  "Automated snapshot restore"),
    ("5", "Defensibility Index",   "GET  /api/v1/dashboard/defensibility-index", ACCENT_BLUE, "Score: 87/100"),
]
colors2 = [ACCENT_RED, ACCENT_BLUE, ACCENT_GRN, ACCENT_GRN, ACCENT_BLUE]

for i, (num, title, endpoint, color, note) in enumerate(steps):
    y = 1.1 + i * 1.1
    box(s, 0.4, y, 0.55, 0.55, fill_color=colors2[i])
    txt(s, num,      0.4,  y,      0.55, 0.55, size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title,    1.1,  y,      3.0,  0.35, size=Pt(15), bold=True, color=WHITE)
    txt(s, endpoint, 1.1,  y+0.3,  5.5,  0.3,  size=Pt(11), color=ACCENT_BLUE)
    txt(s, note,     7.5,  y+0.05, 5.3,  0.45, size=Pt(13), color=LIGHT_GRAY)

slide_number(s, 4)

# ──────────────────────────────────────────────
# SLIDE 5 — Detection Deep Dive (VIGIL)
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, ACCENT_BLUE)

txt(s, "VIGIL — Behavioral Detection Engine", 0.5, 0.2, 12, 0.6, size=Pt(28), bold=True, color=WHITE)

sigs = [
    ("🔴", "Mass Modification",  "5,420 files modified in 60s\nThreshold: 15% — exceeded at 18%\nEntropy Score: 0.89 (encryption confirmed)", ACCENT_RED),
    ("🔴", "VSS Admin Abuse",    "Shadow copy deletion detected\nRecovery prevention attempt\nConfidence: 99%", ACCENT_RED),
    ("🟡", "Lateral Movement",   "Process chain: powershell → cmd\nCross-host credential reuse\nConfidence: 91%", ACCENT_YLW),
    ("🔴", "Entropy Spike",      "Shannon entropy > 0.85\nIndicates active file encryption\nDetected across 3 file types", ACCENT_RED),
]

for i, (icon, title, detail, color) in enumerate(sigs):
    x = 0.4 + (i % 2) * 6.4
    y = 1.1 + (i // 2) * 2.7
    box(s, x, y, 6.0, 2.4, fill_color=RGBColor(0x1A, 0x2B, 0x3C), line_color=color, line_width=Pt(2))
    txt(s, f"{icon}  {title}", x+0.15, y+0.1, 5.7, 0.5, size=Pt(16), bold=True, color=color)
    txt(s, detail, x+0.15, y+0.6, 5.6, 1.6, size=Pt(13), color=LIGHT_GRAY)

txt(s, "ML-based behavioural analysis — detects zero-day ransomware without signature databases",
    0.4, 6.9, 12.5, 0.4, size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 5)

# ──────────────────────────────────────────────
# SLIDE 6 — Grafana & Observability
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, ACCENT_GRN)

txt(s, "Observability — Grafana + Prometheus", 0.5, 0.2, 12, 0.6, size=Pt(28), bold=True, color=WHITE)

txt(s, "http://localhost:7600", 0.5, 0.85, 5, 0.35, size=Pt(13), color=ACCENT_BLUE)

panels = [
    ("Detection Latency",     "Real-time ms counter\nfrom VIGIL events",          ACCENT_BLUE),
    ("Files Affected",        "Rate of file modifications\nper second",             ACCENT_RED),
    ("Containment Time",      "MTTC trend over time\nTarget: < 60 seconds",        ACCENT_YLW),
    ("Defensibility Index",   "Live score gauge\n0–100, higher = better",          ACCENT_GRN),
    ("Kafka Throughput",      "Events/sec through\nthe detection pipeline",         ACCENT_BLUE),
    ("Recovery Progress",     "% data restored\nfrom snapshot",                    ACCENT_GRN),
]

for i, (title, desc, color) in enumerate(panels):
    x = 0.4 + (i % 3) * 4.2
    y = 1.3 + (i // 3) * 2.5
    box(s, x, y, 3.8, 2.1, fill_color=RGBColor(0x1A, 0x2B, 0x3C), line_color=color, line_width=Pt(1.5))
    box(s, x, y, 3.8, 0.45, fill_color=color)
    txt(s, title, x+0.1, y+0.05, 3.6, 0.4, size=Pt(13), bold=True, color=WHITE)
    txt(s, desc,  x+0.1, y+0.55, 3.6, 1.4, size=Pt(12), color=LIGHT_GRAY)

txt(s, "Prometheus scrapes metrics every 15s  •  Grafana provides real-time dashboards  •  Full audit trail in PostgreSQL",
    0.4, 6.9, 12.5, 0.4, size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 6)

# ──────────────────────────────────────────────
# SLIDE 7 — Defensibility Index
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, ACCENT_BLUE)

txt(s, "The Defensibility Index", 0.5, 0.2, 12, 0.6, size=Pt(32), bold=True, color=WHITE)
txt(s, "A single score that answers: How well did you defend?", 0.5, 0.85, 12, 0.4,
    size=Pt(17), color=LIGHT_GRAY)

# Big score
box(s, 4.7, 1.4, 3.9, 2.5, fill_color=RGBColor(0x1A, 0x2B, 0x3C), line_color=ACCENT_BLUE, line_width=Pt(3))
txt(s, "87", 4.7, 1.5, 3.9, 1.5, size=Pt(80), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
txt(s, "out of 100", 4.7, 2.9, 3.9, 0.5, size=Pt(16), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txt(s, "76th percentile vs industry", 4.7, 3.4, 3.9, 0.4, size=Pt(13), color=ACCENT_GRN, align=PP_ALIGN.CENTER)

# Sub scores
sub = [
    ("Detection",    92, ACCENT_BLUE),
    ("Isolation",    85, ACCENT_GRN),
    ("Recovery",     88, ACCENT_GRN),
    ("Immutability", 79, ACCENT_YLW),
]
for i, (label, score, color) in enumerate(sub):
    x = 0.4 + (i % 2) * 4.2
    y = 4.1 + (i // 2) * 1.1
    box(s, x, y, 3.8, 0.85, fill_color=RGBColor(0x1A, 0x2B, 0x3C), line_color=color, line_width=Pt(1.5))
    txt(s, label,        x+0.15, y+0.1,  2.5, 0.4, size=Pt(14), color=WHITE)
    txt(s, f"{score}/100", x+2.8,  y+0.1,  0.9, 0.4, size=Pt(14), bold=True, color=color, align=PP_ALIGN.RIGHT)

txt(s, "Unlike CVE scores — higher is better. Guides security investment decisions.",
    0.4, 6.4, 12.5, 0.5, size=Pt(14), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 7)

# ──────────────────────────────────────────────
# SLIDE 8 — Results Summary
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, ACCENT_GRN)

txt(s, "Demo Results", 0.5, 0.2, 12, 0.6, size=Pt(32), bold=True, color=WHITE)

results = [
    ("⏱️",  "Detection Time",        "8 seconds",   ACCENT_BLUE),
    ("🛡️",  "Attack Phases Caught",   "4 / 4",       ACCENT_GRN),
    ("📊",  "Accuracy",               "95%",         ACCENT_GRN),
    ("💾",  "Data Recovered",         "99.8%",       ACCENT_GRN),
    ("🚫",  "Damage Prevented",       "$2.3M",       ACCENT_YLW),
    ("✅",  "Defensibility Index",    "87 / 100",    ACCENT_BLUE),
]

for i, (icon, label, value, color) in enumerate(results):
    x = 0.4 + (i % 3) * 4.2
    y = 1.1 + (i // 3) * 2.3
    box(s, x, y, 3.8, 2.0, fill_color=RGBColor(0x1A, 0x2B, 0x3C), line_color=color, line_width=Pt(2))
    txt(s, icon,  x+0.15, y+0.1,  3.5, 0.5, size=Pt(22), align=PP_ALIGN.CENTER)
    txt(s, value, x+0.15, y+0.5,  3.5, 0.8, size=Pt(28), bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.15, y+1.3,  3.5, 0.5, size=Pt(12), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "MTTC: 45s  (target: 60s)  •  False positive rate: < 1%  •  Zero manual intervention",
    0.4, 5.8, 12.5, 0.4, size=Pt(14), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 8)

# ──────────────────────────────────────────────
# SLIDE 9 — Key Takeaways
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)

txt(s, "Key Takeaways", 0.5, 0.2, 12, 0.6, size=Pt(32), bold=True, color=WHITE)

takeaways = [
    (ACCENT_BLUE, "Automated Detection",
     "ML-based behavioural analysis catches zero-day ransomware in seconds — no signatures required."),
    (ACCENT_YLW,  "Distributed & Resilient",
     "No single point of failure. Kafka, Redis, and PostgreSQL ensure the pipeline survives partial outages."),
    (ACCENT_GRN,  "Recovery-First Design",
     "Immutable MinIO storage and automated snapshot restore mean data is always recoverable."),
    (ACCENT_RED,  "Defensibility First",
     "The Defensibility Index shifts focus from 'are we vulnerable?' to 'how well did we defend?'"),
]

for i, (color, title, body) in enumerate(takeaways):
    y = 1.1 + i * 1.45
    box(s, 0.4, y, 0.08, 1.1, fill_color=color)
    txt(s, title, 0.65, y,      12, 0.45, size=Pt(17), bold=True, color=color)
    txt(s, body,  0.65, y+0.45, 12, 0.65, size=Pt(14), color=LIGHT_GRAY)

slide_number(s, 9)

# ──────────────────────────────────────────────
# SLIDE 10 — Call to Action
# ──────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, ACCENT_BLUE)

txt(s, "Get Started", 0.5, 0.3, 12, 0.6, size=Pt(32), bold=True, color=WHITE)

actions = [
    ("GitHub",   "github.com/nddmars/DRRA",          "Star the repo, fork it, run it locally"),
    ("API Docs", "localhost:8000/docs",               "Interactive Swagger UI — try every endpoint"),
    ("Grafana",  "localhost:7600",                    "Live metrics dashboard — Prometheus backed"),
    ("Postman",  "DRRA-Postman-Collection.json",      "Import the collection — full demo in one click"),
]

for i, (label, link, desc) in enumerate(actions):
    y = 1.2 + i * 1.35
    box(s, 0.4, y, 2.0, 0.9, fill_color=ACCENT_BLUE)
    txt(s, label, 0.4, y+0.15, 2.0, 0.6, size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, link,  2.6, y,      6.5, 0.45, size=Pt(15), bold=True, color=ACCENT_BLUE)
    txt(s, desc,  2.6, y+0.45, 9.5, 0.4,  size=Pt(13), color=LIGHT_GRAY)

box(s, 1.5, 6.5, 10.3, 0.7, fill_color=ACCENT_BLUE)
txt(s, "\"DRRA shifts ransomware defense from reactive to proactive — measure what matters.\"",
    1.5, 6.5, 10.3, 0.7, size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(s, 10)

# ──────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────
out = r"C:\Users\txvat\git\ransomwaredefense\DRRA_Conference_Demo.pptx"
prs.save(out)
print(f"Saved: {out}")
